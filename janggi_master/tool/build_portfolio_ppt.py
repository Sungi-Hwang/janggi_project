from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "test_tmp" / "portfolio_ppt"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PPT_PATH = OUT_DIR / "janggi_hansu_portfolio_v1.pptx"

HOME_IMG = ROOT / "test_tmp" / "portfolio_home.png"
PUZZLE_IMG = ROOT / "test_tmp" / "portfolio_puzzle_categories.png"
CONTINUE_IMG = ROOT / "test_tmp" / "portfolio_continue.png"
AI_IMG = ROOT / "test_tmp" / "portfolio_ai_setup.png"


BG = RGBColor(247, 242, 234)
PANEL = RGBColor(255, 252, 247)
TEXT = RGBColor(53, 38, 32)
MUTED = RGBColor(108, 92, 82)
ACCENT = RGBColor(118, 79, 53)
ACCENT_LIGHT = RGBColor(230, 213, 193)
BLUE = RGBColor(53, 88, 144)
GREEN = RGBColor(46, 122, 87)
RED = RGBColor(171, 73, 63)


def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(5.9), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Malgun Gothic"
        r2.font.size = Pt(11)
        r2.font.color.rgb = MUTED


def add_body_text(slide, left, top, width, height, lines, font_size=15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.level = 0
        r = p.add_run()
        r.text = line
        r.font.name = "Malgun Gothic"
        r.font.size = Pt(font_size)
        r.font.color.rgb = TEXT


def add_bullets(slide, left, top, width, height, items, font_size=16, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.bullet = True
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = item
        r.font.name = "Malgun Gothic"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color


def add_card(slide, left, top, width, height, title, body, accent=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT_LIGHT

    title_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.18), width - Inches(0.4), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = accent

    body_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.58), width - Inches(0.4), height - Inches(0.72))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for i, line in enumerate(body):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = line
        r.font.name = "Malgun Gothic"
        r.font.size = Pt(13)
        r.font.color.rgb = TEXT


def add_image_card(slide, img_path, left, top, width, height, title):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ACCENT_LIGHT

    slide.shapes.add_picture(str(img_path), left + Inches(0.14), top + Inches(0.14), width=width - Inches(0.28))
    cap = slide.shapes.add_textbox(left + Inches(0.14), top + height - Inches(0.52), width - Inches(0.28), Inches(0.3))
    tf = cap.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Malgun Gothic"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = TEXT


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.55), Inches(6.25), Inches(5.95))
banner.fill.solid()
banner.fill.fore_color.rgb = PANEL
banner.line.color.rgb = ACCENT_LIGHT
slide.shapes.add_picture(str(HOME_IMG), Inches(7.1), Inches(0.55), height=Inches(6.15))
add_title(
    slide,
    "장기 한수",
    "개인 프로젝트 | Flutter · Dart · Fairy-Stockfish FFI · 퍼즐 검증 파이프라인",
)
add_body_text(
    slide,
    Inches(1.0),
    Inches(1.55),
    Inches(5.5),
    Inches(3.7),
    [
        "혼자 장기 실력을 늘리는 훈련형 모바일 앱",
        "AI 대국, 묘수풀이, 이어하기 분석을 하나의 흐름으로 연결",
        "기보 기반 퍼즐 추출과 합법수·체크메이트 검증 로직까지 직접 설계",
    ],
    font_size=18,
)
tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.35), Inches(2.35), Inches(0.5))
tag.fill.solid()
tag.fill.fore_color.rgb = ACCENT
tag.line.color.rgb = ACCENT
tf = slide.shapes.add_textbox(Inches(1.1), Inches(5.45), Inches(2.2), Inches(0.22)).text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Portfolio Project"
r.font.name = "Malgun Gothic"
r.font.size = Pt(12)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 255, 255)


# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "문제 정의", "단순 장기 대국 앱이 아니라, 혼자 복기하고 수읽기를 훈련할 수 있는 앱을 목표로 했습니다.")
add_card(
    slide,
    Inches(0.8),
    Inches(1.5),
    Inches(3.85),
    Inches(4.8),
    "기존 한계",
    [
        "대국 앱은 많지만 실력 향상에 초점을 둔 흐름은 약함",
        "묘수풀이의 정답 판정이 사용자 체감과 어긋나는 경우가 많음",
        "원하는 배치에서 바로 분석·복기하는 기능이 부족함",
    ],
    accent=RED,
)
add_card(
    slide,
    Inches(4.75),
    Inches(1.5),
    Inches(3.85),
    Inches(4.8),
    "해결 방향",
    [
        "AI 대국 + 퍼즐 + 이어하기 분석을 한 앱 안에서 연결",
        "저장된 정답 수순보다 실제 체크메이트 성립 여부를 우선",
        "직접 만든 배치에서 AI 대전 또는 로컬 분석까지 지원",
    ],
    accent=BLUE,
)
add_card(
    slide,
    Inches(8.7),
    Inches(1.5),
    Inches(3.85),
    Inches(4.8),
    "대상 사용자",
    [
        "장기를 좋아하지만 자주 둘 상대가 없는 사용자",
        "실전 복기와 수읽기 훈련을 반복하고 싶은 사용자",
        "퍼즐을 풀고 공유하며 학습하고 싶은 사용자",
    ],
    accent=GREEN,
)


# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "핵심 기능", "사용자 경험은 세 가지 축으로 설계했습니다.")
add_image_card(slide, HOME_IMG, Inches(0.8), Inches(1.55), Inches(3.75), Inches(4.95), "메인 메뉴")
add_image_card(slide, PUZZLE_IMG, Inches(4.8), Inches(1.55), Inches(3.75), Inches(4.95), "묘수풀이 카탈로그")
add_image_card(slide, CONTINUE_IMG, Inches(8.8), Inches(1.55), Inches(3.75), Inches(4.95), "배치 이어하기 / 로컬 분석")


# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "기술 구현", "앱 로직과 엔진, 퍼즐 데이터 파이프라인을 분리해 설계했습니다.")
add_card(
    slide,
    Inches(0.8),
    Inches(1.55),
    Inches(4.0),
    Inches(4.8),
    "앱 구조",
    [
        "Flutter / Dart 기반 모바일 앱",
        "보드 상태와 장기 규칙은 앱 내부 로직으로 직접 관리",
        "퍼즐 진행, 합법수 계산, 장군/체크메이트 판정, undo를 상태 객체에서 통합 처리",
    ],
    accent=ACCENT,
)
add_card(
    slide,
    Inches(4.95),
    Inches(1.55),
    Inches(3.85),
    Inches(4.8),
    "엔진 연동",
    [
        "Fairy-Stockfish를 Dart FFI로 연결",
        "AI 대국, 힌트, 퍼즐 응수 생성에 활용",
        "엔진 결과를 앱 규칙과 다시 대조해 실제 진행 흐름에 반영",
    ],
    accent=BLUE,
)
add_card(
    slide,
    Inches(8.95),
    Inches(1.55),
    Inches(3.6),
    Inches(4.8),
    "데이터/검증",
    [
        "GIB 기보 기반 퍼즐 추출",
        "합법수, 중복 정답, 체크메이트 성립 여부를 기준으로 퍼즐 필터링",
        "현재 strict 퍼즐셋 197개 운영",
    ],
    accent=GREEN,
)


# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "디버깅과 품질 개선", "사용자 피드백을 기반으로 로직과 데이터 둘 다 구조적으로 정리했습니다.")
add_bullets(
    slide,
    Inches(0.95),
    Inches(1.55),
    Inches(5.9),
    Inches(4.9),
    [
        "퍼즐 정답 수순만 강제하던 방식을 수정해, 주어진 수 안에 실제 체크메이트면 성공으로 처리",
        "기보 기반 퍼즐 중 애매하거나 품질이 낮은 문제를 재검증해 카탈로그 정리",
        "궁성 대각선 차/포 이동과 장군 판정 불일치를 수정하고 회귀 테스트 추가",
        "힌트 UI, 퍼즐 진행 표시, 로컬 분석 기능 등 사용자 체감이 큰 부분부터 개선",
    ],
    font_size=18,
)
slide.shapes.add_picture(str(AI_IMG), Inches(7.35), Inches(1.7), height=Inches(4.85))


# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "베타 테스트와 배포 경험", "커뮤니티 모집 → 피드백 수집 → 수정 → 재배포의 사이클을 직접 운영했습니다.")
add_card(
    slide,
    Inches(0.8),
    Inches(1.6),
    Inches(4.0),
    Inches(4.7),
    "피드백 수집",
    [
        "장기 커뮤니티와 오픈채팅을 통해 테스트 유저 모집",
        "퍼즐 오답 판정, 장군 처리, UI 직관성 관련 제보 수집",
        "스크린샷과 재현 과정을 기반으로 버그 분석",
    ],
    accent=RED,
)
add_card(
    slide,
    Inches(4.95),
    Inches(1.6),
    Inches(4.0),
    Inches(4.7),
    "수정 사이클",
    [
        "에뮬레이터·실기기 재현",
        "회귀 테스트 추가",
        "비공개 테스트 트랙 재배포",
        "Play Console 대응 및 릴리스 관리",
    ],
    accent=BLUE,
)
add_card(
    slide,
    Inches(9.1),
    Inches(1.6),
    Inches(3.45),
    Inches(4.7),
    "프로젝트 의미",
    [
        "모델 구현을 넘어 제품 완성까지 경험",
        "데이터, 규칙 로직, 모바일 UX, 배포를 하나의 시스템으로 다룸",
        "실제 사용자 피드백 기반 개선 경험 확보",
    ],
    accent=GREEN,
)


# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title(slide, "회고", "이 프로젝트를 통해 분석과 구현, 검증과 배포를 하나의 흐름으로 다루는 경험을 얻었습니다.")
add_bullets(
    slide,
    Inches(0.95),
    Inches(1.55),
    Inches(11.3),
    Inches(3.7),
    [
        "문제 정의에서 끝나지 않고 실제로 작동하는 구조까지 완성하는 능력을 강화함",
        "사용자 피드백을 데이터와 로직 개선으로 연결하는 경험을 반복함",
        "앱 개발, 엔진 연동, 검증 자동화, 베타 운영을 한 프로젝트 안에서 경험함",
    ],
    font_size=19,
)
footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(5.65), Inches(11.35), Inches(0.85))
footer.fill.solid()
footer.fill.fore_color.rgb = ACCENT
footer.line.color.rgb = ACCENT
box = slide.shapes.add_textbox(Inches(1.25), Inches(5.9), Inches(10.8), Inches(0.3))
tf = box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Data Scientist & Engineer 관점에서 본 개인 프로젝트: 문제정의 → 시스템 설계 → 검증 → 서비스 적용"
r.font.name = "Malgun Gothic"
r.font.size = Pt(17)
r.font.bold = True
r.font.color.rgb = RGBColor(255, 255, 255)


prs.save(PPT_PATH)
print(PPT_PATH)
